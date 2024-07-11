import sys
from PyQt6.QtWidgets import QApplication, QMainWindow, QVBoxLayout, QWidget, QPushButton, QFileDialog, QLabel
from PyQt6.QtGui import QPixmap
from PyQt6.QtCore import Qt
import collections 
import collections.abc
from pptx import Presentation
from PIL import Image
import io

class PowerPointViewer(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PyQt6 PowerPoint Viewer")
        self.setGeometry(100, 100, 800, 600)

        self.central_widget = QWidget()
        self.setCentralWidget(self.central_widget)
        self.layout = QVBoxLayout(self.central_widget)

        self.open_button = QPushButton("Open PowerPoint")
        self.open_button.clicked.connect(self.open_file)
        self.layout.addWidget(self.open_button)

        self.image_label = QLabel()
        self.image_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.layout.addWidget(self.image_label)

        self.current_slide = 0
        self.slides = []

        self.next_button = QPushButton("Next Slide")
        self.next_button.clicked.connect(self.next_slide)
        self.layout.addWidget(self.next_button)

        self.prev_button = QPushButton("Previous Slide")
        self.prev_button.clicked.connect(self.prev_slide)
        self.layout.addWidget(self.prev_button)

    def open_file(self):
        file_name, _ = QFileDialog.getOpenFileName(self, "Open PowerPoint", "", "PowerPoint Files (*.pptx)")
        if file_name:
            self.load_presentation(file_name)

    def load_presentation(self, file_name):
        prs = Presentation(file_name)
        self.slides = []
        for slide in prs.slides:
            img = self.slide_to_image(slide)
            self.slides.append(img)
        self.current_slide = 0
        self.show_current_slide()

    def slide_to_image(self, slide):
        img_stream = io.BytesIO()
        slide.shapes.add_picture = lambda *args, **kwargs: None  # Workaround for unsupported images
        slide.save(img_stream, 'PNG')
        img_stream.seek(0)
        return Image.open(img_stream)

    def show_current_slide(self):
        if self.slides:
            img = self.slides[self.current_slide]
            qimage = self.pil_image_to_qpixmap(img)
            self.image_label.setPixmap(qimage)

    def pil_image_to_qpixmap(self, pil_image):
        bytes_img = io.BytesIO()
        pil_image.save(bytes_img, format='PNG')
        qimg = QPixmap()
        qimg.loadFromData(bytes_img.getvalue())
        return qimg

    def next_slide(self):
        if self.current_slide < len(self.slides) - 1:
            self.current_slide += 1
            self.show_current_slide()

    def prev_slide(self):
        if self.current_slide > 0:
            self.current_slide -= 1
            self.show_current_slide()

if __name__ == "__main__":
    app = QApplication(sys.argv)
    viewer = PowerPointViewer()
    viewer.show()
    sys.exit(app.exec())